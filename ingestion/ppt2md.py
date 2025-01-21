import base64
import logging
from openai import OpenAI
from pptx import Presentation
from pathlib import Path
from markitdown import MarkItDown
import sys
import os
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
from utils.init import key
# Initialize logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Initialize the LLM client (replace `key` with your actual API key)

client = OpenAI(api_key=key)

# Function to encode the image to Base64
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

# Function to extract images from all slides of a PPTX file
def extract_images_from_pptx(pptx_path, output_folder):
    prs = Presentation(pptx_path)
    image_data = []

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for slide_idx, slide in enumerate(prs.slides):
        logging.info(f"Processing slide {slide_idx + 1}")
        slide_info = {"slide_number": slide_idx + 1, "images": []}
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # 13 = Picture
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_path = os.path.join(output_folder, f"slide_{slide_idx + 1}_img_{shape_idx + 1}.{image_ext}")
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)
                slide_info["images"].append(image_path)
                logging.info(f"Extracted image: {image_path}")
        image_data.append(slide_info)

    return image_data

# Function to describe images using GPT-4 Vision
def describe_image_with_gpt4(image_path):
    try:
        base64_image = encode_image(image_path)
        logging.info(f"Encoding image: {image_path}")

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Describe this image in detail."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                    ],
                }
            ]
        )
        description = response.choices[0].message.content.strip()
        logging.info(f"Description for {image_path}: {description}")
        return description
    except Exception as e:
        logging.error(f"Error describing {image_path}: {e}")
        return "Description unavailable."

# Function to generate Markdown from PPTX using MarkItDown and enhance with descriptions
def generate_markdown_with_markitdown(pptx_path, image_data):
    # Convert PPTX to base Markdown using MarkItDown
    md = MarkItDown(llm_client=client, llm_model="gpt-4o-mini")
    result = md.convert(pptx_path)
    markdown_lines = result.text_content.splitlines()

    # Enhance Markdown with image descriptions
    enhanced_markdown = []
    slide_idx = 0  # Track the current slide index in image_data

    for line in markdown_lines:
        enhanced_markdown.append(line)

        # Check if this line marks the start of a slide section
        if line.strip().startswith("<!-- Slide number:"):
            slide_number = int(line.strip().split(":")[1].strip().strip("-->"))

            # Add descriptions for the current slide
            if slide_idx < len(image_data) and image_data[slide_idx]["slide_number"] == slide_number:
                for image_path in image_data[slide_idx]["images"]:
                    description = describe_image_with_gpt4(image_path)
                    enhanced_markdown.append(f"- **Image**: {Path(image_path).name}")
                    enhanced_markdown.append(f"  - Description: {description}")
                slide_idx += 1  # Move to the next slide in image_data

    return "\n".join(enhanced_markdown)

# Main workflow
def convert_pptx_to_markdown_with_markitdown(pptx_path, output_folder, markdown_output):
    # Extract images from all slides of the PPTX file
    image_data = extract_images_from_pptx(pptx_path, output_folder)

    # Generate Markdown with MarkItDown and enhance with descriptions
    markdown_content = generate_markdown_with_markitdown(pptx_path, image_data)

    # Save the Markdown to a file
    with open(markdown_output, "w") as md_file:
        md_file.write(markdown_content)

    logging.info(f"Markdown with enhanced descriptions saved to {markdown_output}")

# Example usage
if __name__ == "__main__":
    pptx_path = "data/taxgpt/raw/MIC_3e_Ch11.pptx"  # Input PPTX file
    output_folder = "data/taxgpt/processed/images/MIC_3e_Ch11"  # Folder for extracted images
    markdown_output = "data/taxgpt/processed/MIC_3e_Ch11.md"  # Output Markdown file

    convert_pptx_to_markdown_with_markitdown(
        pptx_path=pptx_path,
        output_folder=output_folder,
        markdown_output=markdown_output
    )
